"""
PPT生成器 - 将提取的帧和分析结果生成PowerPoint
支持多种模板风格
"""

import os
import io
import tempfile

from PIL import Image

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor as PPTXColor
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False


# 预设模板配色
TEMPLATES = {
    'modern': {
        'name': '现代简洁',
        'slide_bg': (255, 255, 255),
        'title_color': (33, 37, 41),
        'text_color': (60, 60, 60),
        'accent_color': (66, 133, 244),
        'slide_bg_image': False,
    },
    'dark': {
        'name': '深色商务',
        'slide_bg': (30, 33, 45),
        'title_color': (255, 255, 255),
        'text_color': (200, 200, 210),
        'accent_color': (96, 165, 250),
        'slide_bg_image': False,
    },
    'gradient': {
        'name': '渐变科技',
        'slide_bg': (248, 250, 252),
        'title_color': (15, 23, 42),
        'text_color': (71, 85, 105),
        'accent_color': (139, 92, 246),
        'slide_bg_image': False,
    },
    'green': {
        'name': '清新绿',
        'slide_bg': (255, 255, 255),
        'title_color': (20, 83, 45),
        'text_color': (40, 50, 60),
        'accent_color': (34, 197, 94),
        'slide_bg_image': False,
    },
}


class PPTGenerator:
    def __init__(self, template='modern', title='视频要点提取', subtitle='自动生成的演示文稿'):
        if not _HAS_PPTX:
            raise RuntimeError("请先安装 python-pptx: pip install python-pptx")
        self.template_name = template
        self.template = TEMPLATES.get(template, TEMPLATES['modern'])
        self.presentation = Presentation()
        # 16:9 宽屏
        self.presentation.slide_width = Inches(13.333)
        self.presentation.slide_height = Inches(7.5)
        self.title_text = title
        self.subtitle_text = subtitle

    def _rgb(self, rgb_tuple):
        return PPTXColor(*rgb_tuple)

    def _set_background(self, slide, color):
        """给幻灯片设置背景色"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._rgb(color)

    def _add_accent_bar(self, slide, x, y, width, height):
        """在幻灯片左上角添加一个装饰色条"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(self.template['accent_color'])
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def add_title_slide(self):
        """添加封面页"""
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])  # 空白布局
        self._set_background(slide, self.template['slide_bg'])

        # 左侧装饰色条
        self._add_accent_bar(slide, Inches(0.5), Inches(2.8), Inches(1.5), Inches(0.08))

        # 标题
        left = Inches(1.5)
        top = Inches(2.2)
        width = Inches(10.5)
        height = Inches(1.2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = self.title_text
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self._rgb(self.template['title_color'])
        p.alignment = PP_ALIGN.LEFT

        # 副标题
        sub = txBox.text_frame.add_paragraph()
        sub.text = self.subtitle_text
        sub.font.size = Pt(24)
        sub.font.color.rgb = self._rgb(self.template['accent_color'])
        sub.space_before = Pt(12)

        # 底部信息
        bottom_box = slide.shapes.add_textbox(Inches(1.5), Inches(6.2), Inches(10), Inches(0.6))
        bp = bottom_box.text_frame.paragraphs[0]
        bp.text = "VideoToPPT · 自动生成"
        bp.font.size = Pt(14)
        bp.font.color.rgb = self._rgb(self.template['text_color'])

    def add_content_slide(self, img_rgb, title, bullet_points=None,
                          image_layout='right', include_original_image=True):
        """
        添加一张内容页

        参数:
            img_rgb: numpy.ndarray (H, W, 3) RGB
            title: 幻灯片标题
            bullet_points: list[str] 要点列表
            image_layout: 'right' | 'left' | 'full' | 'top'  图片布局
            include_original_image: 是否包含原始截图
        """
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        self._set_background(slide, self.template['slide_bg'])

        # 顶部色条装饰
        self._add_accent_bar(slide, Inches(0.4), Inches(0.4), Inches(0.6), Inches(0.12))

        # 标题
        title_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.4), Inches(11.5), Inches(1.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        display_title = title if title else "无标题"
        if len(display_title) > 80:
            display_title = display_title[:80] + "..."
        p.text = display_title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self._rgb(self.template['title_color'])

        # 根据布局放置图片和文本
        if image_layout == 'full':
            # 图片占全屏，文本在底部
            img_path = self._temp_image_path(img_rgb)
            left = Inches(0.4)
            top = Inches(1.4)
            width = Inches(12.5)
            height = Inches(4.5)
            slide.shapes.add_picture(img_path, left, top, width=width, height=height)

            if bullet_points:
                text_box = slide.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(11.5), Inches(1.3))
                self._add_bullet_text(text_box, bullet_points)
            self._cleanup_temp(img_path)

        elif image_layout == 'top':
            # 图片上半部分，文本下半部分
            img_path = self._temp_image_path(img_rgb)
            left = Inches(0.4)
            top = Inches(1.4)
            width = Inches(12.5)
            height = Inches(3.2)
            slide.shapes.add_picture(img_path, left, top, width=width, height=height)
            self._cleanup_temp(img_path)

            if bullet_points:
                text_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.7), Inches(11.5), Inches(2.5))
                self._add_bullet_text(text_box, bullet_points)

        else:  # right or left (图文并排)
            if image_layout == 'left':
                img_left = Inches(0.5)
                text_left = Inches(6.5)
            else:
                img_left = Inches(7.0)
                text_left = Inches(1.0)

            # 图片
            if include_original_image:
                img_path = self._temp_image_path(img_rgb)
                top = Inches(1.6)
                width = Inches(5.8)
                height = Inches(5.2)
                pic = slide.shapes.add_picture(img_path, img_left, top, width=width, height=height)
                # 加阴影/边框效果 - 通过添加一个背景矩形
                shadow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, img_left + Inches(0.05), top + Inches(0.08), width, height)
                shadow.fill.solid()
                shadow.fill.fore_color.rgb = self._rgb(self.template['accent_color'])
                shadow.line.fill.background()
                # 移动到图片后面
                shadow.shadow.inherit = False
                slide.shapes._spTree.remove(shadow._element)
                slide.shapes._spTree.insert(2, shadow._element)
                self._cleanup_temp(img_path)

            # 文本
            text_box = slide.shapes.add_textbox(text_left, Inches(1.8), Inches(5.5), Inches(5.0))
            self._add_bullet_text(text_box, bullet_points or [])

    def _add_bullet_text(self, text_box, bullet_points):
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP

        if not bullet_points:
            p = tf.paragraphs[0]
            p.text = "（无文本内容）"
            p.font.size = Pt(18)
            p.font.italic = True
            p.font.color.rgb = self._rgb(self.template['text_color'])
            return

        for i, bp in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            display = bp.strip()
            if len(display) > 120:
                display = display[:120] + "..."
            p.text = "• " + display
            p.font.size = Pt(20)
            p.font.color.rgb = self._rgb(self.template['text_color'])
            p.space_after = Pt(8)

    def _temp_image_path(self, img_rgb):
        """将numpy图像保存为临时文件并返回路径"""
        pil = Image.fromarray(img_rgb)
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
        tmp.close()
        pil.save(tmp.name, 'PNG')
        self._last_temp = tmp.name
        return tmp.name

    def _cleanup_temp(self, path):
        try:
            if os.path.exists(path):
                os.unlink(path)
        except Exception:
            pass

    def save(self, output_path):
        """保存PPT到文件"""
        self.presentation.save(output_path)
        return output_path


def build_presentation(frames_with_analysis, output_path,
                       template='modern',
                       title='视频要点',
                       subtitle='由 VideoToPPT 自动生成',
                       image_layout='right',
                       include_ocr=True,
                       progress_callback=None):
    """
    从分析结果构建完整PPT

    参数:
        frames_with_analysis: list[dict] - 每帧包含 'image'(RGB), 'analysis'(dict), 'frame_index', 'timestamp'
        output_path: 输出文件路径
        template: 模板名
        image_layout: 'right' | 'left' | 'full' | 'top'
        progress_callback: 进度回调 (percent, message)
    """
    gen = PPTGenerator(template=template, title=title, subtitle=subtitle)
    gen.add_title_slide()

    total = len(frames_with_analysis)
    for i, item in enumerate(frames_with_analysis):
        img = item.get('image')
        analysis = item.get('analysis', {})
        title_txt = analysis.get('title') or analysis.get('auto_title') or f"幻灯片 {i + 1}"

        bullets = []
        if include_ocr and analysis.get('bullet_points'):
            bullets = analysis['bullet_points']
        elif include_ocr and analysis.get('content'):
            bullets = [b for b in analysis['content'].split('\n') if b.strip()]

        gen.add_content_slide(
            img_rgb=img,
            title=title_txt,
            bullet_points=bullets,
            image_layout=image_layout,
            include_original_image=True,
        )
        if progress_callback:
            percent = int(50 + 50.0 * (i + 1) / max(total, 1))
            progress_callback(percent, f"生成第 {i + 1}/{total} 页...")

    gen.save(output_path)
    return output_path
